import configparser
import pandas
import pptx
import pytest
import pronto.pronto

from contextlib import nullcontext as does_not_raise

@pytest.mark.parametrize(
    "input, exception, want",
    [
        (-1, does_not_raise(), "-1 mut/Mb; Not available\n"),
        (2.5, does_not_raise(), "2.5 mut/Mb; Lav\n"),
        (7, does_not_raise(), "7 mut/Mb; Intermediær\n"),
        (23, does_not_raise(), "23 mut/Mb; Høy\n"),
    ]
)
def test_get_tmb_string(input, exception, want):
    with exception:
        assert pronto.pronto.get_tmb_string(input) == want

@pytest.mark.parametrize(
    "inputs, exception, want",
    [
        (
            (
                True,
                'pronto/tests/data/ous',
                '250101_NDX012345_RUO_0001_01234ABCDE',
                'IPH0001-D01-R01-A01',
                'test.txt'
            ),
            does_not_raise(),
            'pronto/tests/data/ous/250101_NDX012345_RUO_0001_01234ABCDE_TSO_500_LocalApp_postprocessing_results/IPH0001-D01-R01-A01/test.txt'
        ),
        (
            (
                True,
                'pronto/tests/data/hus/',
                '250101_NDX012345_RUO_0001_01234ABCDE',
                'IPH0001-D01-R01-A01',
                'test*.txt'
            ),
            does_not_raise(),
            'pronto/tests/data/hus/250101_NDX012345_RUO_0001_01234ABCDE/250101_NDX012345_RUO_0001_01234ABCDE_ppr/IPH0001-D01-R01-A01/test.txt'
        ),
        (
            (
                True,
                'pronto/tests/data/hus/',
                '250101_NDX012345_RUO_0001_01234ABCDE',
                'IPH0001-D01-R01-A01',
                'tes*.txt'
            ),
            does_not_raise(),
            'pronto/tests/data/hus/250101_NDX012345_RUO_0001_01234ABCDE/250101_NDX012345_RUO_0001_01234ABCDE_ppr/IPH0001-D01-R01-A01/test.txt'
        ),
        (
            (
                True,
                'pronto/tests/data/hus/',
                '250101_NDX012345_RUO_0001_01234ABCDE',
                'IPH0001-D01-R01-A01',
            ),
            does_not_raise(),
            'pronto/tests/data/hus/250101_NDX012345_RUO_0001_01234ABCDE/250101_NDX012345_RUO_0001_01234ABCDE_ppr/IPH0001-D01-R01-A01'
        ),
        (
            (
                True,
                'pronto/tests/data/none',
                'existent',
                'test.txt'
            ),
            pytest.raises(ValueError),
            None
        ),
        (
            (
                True,
                'pronto/tests/data/*',
                '250101_NDX012345_RUO_0001_01234ABCDE',
                'IPH0001-D01-R01-A01',
                'test.txt'
            ),
            pytest.raises(ValueError),
            None
        ),
        (
            (
                False,
                'pronto/tests/data/*',
                '250101_NDX012345_RUO_0001_01234ABCDE',
                'IPH0001-D01-R01-A01',
                'test.txt'
            ),
            does_not_raise(),
            None
        ),
    ]
)
def test_glob_tsoppi_file(inputs, exception, want):
    with exception:
        assert pronto.pronto.glob_tsoppi_file(*inputs) == want

@pytest.mark.parametrize(
    "inputs, exception, want",
    [
        (
            (
                pandas.DataFrame({
                    "one": [1, 2],
                    "two": [3, 4],
                    "three": [5, 6],
                    "four": [7, 8],
                }),
                ["one", "two", "three", "four"],
            ),
            does_not_raise(),
            pandas.DataFrame({
                "one": [1, 2],
                "two": [3, 4],
                "three": [5, 6],
                "four": [7, 8],
            }),
        ),
        (
            (
                pandas.DataFrame({
                    "one": [1, 2],
                    "two": [3, 4],
                    "four": [7, 8],
                }),
                ["one", "two", "three", "four"],
            ),
            does_not_raise(),
            pandas.DataFrame({
                "one": [1, 2],
                "two": [3, 4],
                "three": [' ', ' '],
                "four": [7, 8],
            }),
        ),
        (
            (
                pandas.DataFrame({
                    "one": [1, 2],
                    "two": [3, 4],
                    "three": [5, 6],
                    "four": [7, 8],
                }),
                ["two", "three", "four"],
            ),
            does_not_raise(),
            pandas.DataFrame({
                "two": [3, 4],
                "three": [5, 6],
                "four": [7, 8],
                "one": [1, 2],
            }),
        ),
        (
            (
                pandas.DataFrame({
                    "one": [1, 2],
                    "two": [3, 4],
                    "four": [7, 8],
                    "five": [9, 10],
                }),
                ["one", "two", "three", "four"],
            ),
            does_not_raise(),
            pandas.DataFrame({
                "one": [1, 2],
                "two": [3, 4],
                "three": [' ', ' '],
                "four": [7, 8],
                "five": [9, 10],
            }),
        ),
    ]
)
def test_normalize_column_index(inputs, exception, want):
    with exception:
        get = pronto.pronto.normalize_column_index(*inputs)
        assert want.equals(get)

@pytest.mark.parametrize(
    "inputs, exception, want",
    [
        (
            (
                pandas.DataFrame({
                    "one": [1, 2],
                    "two": [3.333, 4.444],
                }),
                "two",
            ),
            does_not_raise(),
            pandas.DataFrame({
                "one": [1, 2],
                "two": ["3.33", "4.44"],
            }),
        ),
        (
            (
                pandas.DataFrame({
                    "one": [1, 2],
                    "two": ['21.0%', '0.5%'],
                }),
                "two",
            ),
            does_not_raise(),
            pandas.DataFrame({
                "one": [1, 2],
                "two": ["21.0%", "0.5%"],
            }),
        ),
        (
            (
                pandas.DataFrame({
                    "one": [1, 2],
                }),
                "two",
            ),
            does_not_raise(),
            pandas.DataFrame({
                "one": [1, 2],
            }),
        ),
    ]
)
def test_set_column_to_2_decimals(inputs, exception, want):
    with exception:
        get = pronto.pronto.set_column_to_2_decimals(*inputs)
        assert want.equals(get)

def list_of_lists_equal(list1, list2):
    if len(list1) != len(list2):
        return False
    for sublist1, sublist2 in zip(list1, list2):
        print(sublist1, sublist2)
        if sublist1 != sublist2:
            return False
    return True

@pytest.mark.parametrize(
    "inputs, exception, want",
    [
        (
            (
                pandas.DataFrame({
                    "one": [1, 2],
                    "two": [3, 4],
                }),
                0,
                3,
            ),
            does_not_raise(),
            [
                ["one", "two"],
                [1, 3],
                [2, 4],
            ],
        ),
        (
            (
                pandas.DataFrame({
                    "one": [1, 2, 3 ,4],
                    "two": [5, 6, 7, 8],
                }),
                1,
                2,
            ),
            does_not_raise(),
            [
                ["one", "two"],
                [3, 7],
                [4, 8],
            ],
        ),
        (
            (
                pandas.DataFrame({
                    "one": [1, 2, 3],
                    "two": [5, 6, 7],
                }),
                1,
                2,
            ),
            does_not_raise(),
            [
                ["one", "two"],
                [3, 7],
            ],
        ),
        (
            (
                pandas.DataFrame({
                    "one": [1, 2, 3],
                    "two": [5, 6, 7],
                }),
                2,
                2,
            ),
            does_not_raise(),
            [],
        ),
    ]
)
def test_get_slide_table_data(inputs, exception, want):
    with exception:
        get = pronto.pronto.get_slide_table_data(*inputs)
        assert list_of_lists_equal(get, want)

def check_shape(shape, want_left, want_top, want_width, want_height):
    assert shape.left == pptx.util.Inches(want_left)
    assert shape.top == pptx.util.Inches(want_top)
    assert shape.width == pptx.util.Inches(want_width)
    assert shape.height == pptx.util.Inches(want_height)

def check_paragraph(paragraph, want_text, want_font_size, want_bold, want_alignment):
    assert paragraph.text == want_text
    assert paragraph.font.size.pt == want_font_size
    assert paragraph.font.bold == want_bold
    assert paragraph.alignment == want_alignment

@pytest.mark.parametrize(
    "inputs, exception, want_shape, want_paragraph",
    [
        (
            (
                'Test',
                0.5,
                0.5,
                4,
                1,
                12,
                True,
                0,
                3,
                4,
            ),
            does_not_raise(),
            (
                0.5,
                0.5,
                4,
                1,
            ),
            (
                'Test (N=4, Page 1/3)',
                12.0,
                True,
                pptx.enum.text.PP_ALIGN.CENTER,
            )
        ),
        (
            (
                'Test',
                0.5,
                0.5,
                4,
                1,
                12,
                True,
                0,
                1,
                4,
            ),
            does_not_raise(),
            (
                0.5,
                0.5,
                4,
                1,
            ),
            (
                'Test (N=4)',
                12.0,
                True,
                pptx.enum.text.PP_ALIGN.CENTER,
            )
        ),
        (
            (
                'Test',
                0.5,
                0.5,
                4,
                1,
                12,
                False,
                0,
                3,
                4,
            ),
            does_not_raise(),
            (
                0.5,
                0.5,
                4,
                1,
            ),
            (
                'Test',
                12.0,
                True,
                pptx.enum.text.PP_ALIGN.CENTER,
            )
        ),
    ]
)
def test_add_table_name(inputs, exception, want_shape, want_paragraph):
    with exception:
        shapes = pptx.Presentation().slides.add_slide(pptx.Presentation().slide_layouts[6]).shapes
        pronto.pronto.add_table_name(shapes, *inputs)
        check_shape(shapes[0], *want_shape)
        check_paragraph(shapes[0].text_frame.paragraphs[0], *want_paragraph)

@pytest.mark.parametrize(
    "inputs, exception, want",
    [
        (
            (
                {
                    'FILTER0-1': {
                        'filter_column': 'col1,col2',
                        'key_word': 'keyword1',
                        'columns': 'col1,col2',
                        'filter_column_add': 'col1',
                        'output_table': 'output_table',
                        'min_depth_tumor_dna': '0',
                    }
                },
                "output_dir"
            ),
            does_not_raise(),
            [
                    {
                        'filter_column': 'col1,col2',
                        'key_word': 'keyword1',
                        'columns': 'col1,col2',
                        'filter_column_add': 'col1',
                        'output_table': 'output_table',
                        'min_depth_tumor_dna': '0',
                        'filter_columns': ['col1', 'col2'],
                        'pre_table_output_path': 'output_dir/output_table_pre.txt',
                        'table_output_path': 'output_dir/output_table.txt',
                    }
            ],
        ),
        (
            (
                {
                    'FILTER0-1': {
                        'filter_column': 'col1,col2',
                        'key_word': 'keyword1',
                        'columns': 'col1,col2',
                        'output_table': 'output_table',
                        'min_depth_tumor_dna': '0',
                    }
                },
                "output_dir"
            ),
            does_not_raise(),
            [
                    {
                        'filter_column': 'col1,col2',
                        'key_word': 'keyword1',
                        'columns': 'col1,col2',
                        'filter_column_add': None,
                        'output_table': 'output_table',
                        'min_depth_tumor_dna': '0',
                        'filter_columns': ['col1', 'col2'],
                        'pre_table_output_path': 'output_dir/output_table_pre.txt',
                        'table_output_path': 'output_dir/output_table.txt',
                    }
            ],
        ),
        (
            (
                {
                    'FILTER0-1': {
                        'filter_column': 'col1,col2',
                        'key_word': 'keyword1',
                        'columns': 'col1,col2',
                        'output_table': 'output_table',
                        'filter_column_add': 'col1',
                    }
                },
                "output_dir"
            ),
            does_not_raise(),
            [
                    {
                        'filter_column': 'col1,col2',
                        'key_word': 'keyword1',
                        'columns': 'col1,col2',
                        'filter_column_add': 'col1',
                        'output_table': 'output_table',
                        'min_depth_tumor_dna': None,
                        'filter_columns': ['col1', 'col2'],
                        'pre_table_output_path': 'output_dir/output_table_pre.txt',
                        'table_output_path': 'output_dir/output_table.txt',
                    }
            ],
        ),
        (
            (
                {
                    'FILTER0-1': {
                        'filter_column': 'col1',
                        'key_word': 'keyword1',
                        'columns': 'col1,col2',
                        'output_table': 'output_table',
                        'filter_column_add': 'col1',
                        'min_depth_tumor_dna': '0',
                    }
                },
                "output_dir"
            ),
            does_not_raise(),
            [
                    {
                        'filter_column': 'col1',
                        'key_word': 'keyword1',
                        'columns': 'col1,col2',
                        'filter_column_add': 'col1',
                        'output_table': 'output_table',
                        'min_depth_tumor_dna': '0',
                        'filter_columns': ['col1'],
                        'pre_table_output_path': 'output_dir/output_table_pre.txt',
                        'table_output_path': 'output_dir/output_table.txt',
                    }
            ],
        ),
    ]
)
def test_parse_topfilter(inputs, exception, want):
    with exception:
        cfg = configparser.ConfigParser()
        cfg.read_dict(inputs[0])
        topfilters = pronto.pronto.parse_topfilter(cfg, inputs[1])
        assert topfilters == want

@pytest.mark.parametrize(
    "inputs, exception, want",
    [
        (
            (
                pandas.DataFrame({
                    "filter_column": ["keyword1", "keyword2", "keyword1,keyword2", "other"],
                    "IGV_QC": ["OK", "Not OK", "OK", "OK"],
                    "Class_judgement": ["exclude", "exclude", "include", "include"],
                    "SampleID": ["sample1", "sample1", "sample1", "sample1"],
                }),
                'sample1',
                'filter_column',
                'keyword1',
            ),
            does_not_raise(),
            pandas.DataFrame({
                "filter_column": ["keyword1", "keyword1,keyword2"],
                "IGV_QC": ["OK", "OK"],
                "Class_judgement": ["exclude", "include"],
                "SampleID": ["sample1", "sample1"],
            }),
        ),
        (
            (
                pandas.DataFrame({
                    "filter_column": ["keyword1", "keyword2", "keyword1,keyword2", "other"],
                    "IGV_QC": ["OK", "Not OK", "OK", "OK"],
                    "Class_judgement": ["exclude", "exclude", "include", "include"],
                    "SampleID": ["sample1", "sample1", "sample2", "sample1"],
                }),
                'sample1',
                'filter_column',
                'keyword1',
            ),
            does_not_raise(),
            pandas.DataFrame({
                "filter_column": ["keyword1"],
                "IGV_QC": ["OK"],
                "Class_judgement": ["exclude"],
                "SampleID": ["sample1"],
            }),
        ),
        (
            (
                pandas.DataFrame({
                    "filter_column": ["keyword1", "keyword2", "keyword1,keyword2", "other"],
                    "IGV_QC": ["OK", "Not OK", "OK", "OK"],
                    "Class_judgement": ["exclude", "exclude", "include", "include"],
                    "SampleID": ["sample1", "sample1", "sample1", "sample1"],
                }),
                'sample1',
                'filter_column',
                'keyword1,keyword2',
            ),
            does_not_raise(),
            pandas.DataFrame({
                "filter_column": ["keyword1", "keyword2", "keyword1,keyword2"],
                "IGV_QC": ["OK", "Not OK", "OK"],
                "Class_judgement": ["exclude", "exclude", "include"],
                "SampleID": ["sample1", "sample1", "sample1"],
            }),
        ),
        (
            (
                pandas.DataFrame({
                    "filter_column": ["keyword1", "keyword2", "keyword1,keyword2", "other"],
                    "IGV_QC": ["OK", "Not OK", "OK", "OK"],
                    "Class_judgement": ["exclude", "exclude", "include", "include"],
                    "SampleID": ["sample1", "sample1", "sample1", "sample1"],
                }),
                'sample1',
                'filter_column',
                '!keyword1',
            ),
            does_not_raise(),
            pandas.DataFrame({
                "filter_column": ["keyword2", "other"],
                "IGV_QC": ["Not OK", "OK"],
                "Class_judgement": ["exclude", "include"],
                "SampleID": ["sample1", "sample1"],
            }),
        ),
        (
            (
                pandas.DataFrame({
                    "filter_column": ["keyword1", "keyword2", "keyword1,keyword2", "other"],
                    "IGV_QC": ["OK", "Not OK", "OK", "OK"],
                    "Class_judgement": ["exclude", "exclude", "include", "include"],
                    "SampleID": ["sample1", "sample1", "sample1", "sample1"],
                }),
                'sample1',
                'filter_column',
                '!keyword1 && !keyword2',
            ),
            does_not_raise(),
            pandas.DataFrame({
                "filter_column": ["other"],
                "IGV_QC": ["OK"],
                "Class_judgement": ["include"],
                "SampleID": ["sample1"],
            }),
        ),
        (
            (
                pandas.DataFrame({
                    "filter_column": ["keyword1", "keyword2", "keyword1,keyword2", "other"],
                    "IGV_QC": ["OK", "Not OK", "OK", "OK"],
                    "SampleID": ["sample1", "sample1", "sample1", "sample1"],
                }),
                'sample1',
                'filter_column',
                'keyword1',
            ),
            pytest.raises(ValueError),
            pandas.DataFrame(),
        ),
        (
            (
                pandas.DataFrame({
                    "filter_column": ["keyword1", "keyword2", "keyword1,keyword2", "other"],
                    "IGV_QC": ["OK", "Not OK", "Not OK", "OK"],
                    "Class_judgement": ["exclude", "exclude", "include", "include"],
                    "SampleID": ["sample1", "sample1", "sample1", "sample1"],
                }),
                'sample1',
                'filter_column',
                'keyword1',
            ),
            pytest.raises(ValueError),
            pandas.DataFrame(),
        ),
    ]
)
def test_filter_small_variant_data(inputs, exception, want):
    with exception:
        data = pronto.pronto.filter_small_variant_data(inputs[0], inputs[1], inputs[2], inputs[3])
        data = data.reset_index(drop=True)
        assert data.equals(want)